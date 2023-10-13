import cv2
#import numpy as np
#import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import fitz

def convertPdf2Image() :
  pdf_path = input("Digite o caminho para o arquivo: ")

  pdf_document = fitz.open(pdf_path)
  metadata = pdf_document.metadata

  len_of_document = len(pdf_document)

  for page_number in range(len_of_document):
      page = pdf_document.load_page(page_number)
      
      image = page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
      
      image.save(f'page_{page_number + 1}.png')
      
  pdf_document.close()
  return len_of_document, metadata['title']

def main() :

  index = 1

  prs = Presentation()
  prs.slide_width = Inches(10.5)
  prs.slide_height = Inches(5.85)
  left = top = Inches(0)

  len_of_document, title = convertPdf2Image()
  for i in range(1, len_of_document + 1) :

    path = f"page_{i}.png"

    #print(path)
    img = cv2.imread(path)

    img_gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

    #limiar, img_binary = cv2.threshold(img_gray, 150, 255, cv2.THRESH_BINARY)
    
    #kernel = np.ones((5, 5), np.uint8)
    #img_cleared = cv2.morphologyEx(img_binary, cv2.MORPH_OPEN, kernel)

    img_canny = cv2.Canny(img_gray, 50, 200)
    contourns, _ = cv2.findContours(img_canny, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    #cnt = contourns[4]
    area = cv2.contourArea(contourns[1])
    #per = cv2.arcLength(contourns[4], True)
    #print(len(contourns))
    # print(area)
    # print(per)

    for j in range(len(contourns), 0, -1) :
      #print(j)
      area = cv2.contourArea(contourns[j-1])
      #print(f"Contorno{j}: ")

      min_area = 333410
      
      if area >= min_area :
        #print(area)
        x, y, w, h = cv2.boundingRect(contourns[j-1])
        cropped_img = img[y:y+h, x:x+w]
        #print(index)
        img_name = str(index)+".jpg"
        cv2.imwrite(img_name, cropped_img)
        #cv2.rectangle(img, (x, y), (x+w, y+h), (0, 255, 0), 10)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(img_name, left, top, prs.slide_width, prs.slide_height)
        index += 1
    #cv2.drawContours(img, [contourns[4]], -1, (0, 255, 0), 10)
    # for i in range(len(contourns)) :
    #   print(contourns)
    #   cv2.drawContours(img, contourns, i, (0, 255, 0), 3)
    #   cv2.imshow("Imagem Contornada: ", img)
    #   cv2.waitKey(0)

    #cv2.imshow("Imagem Cinza: ", img_gray)
    #cv2.imshow("Imagem Binarizada: ", img_binary)

    #cv2.waitKey(0)
    #cv2.destroyAllWindows()

    #plt.imshow(img_canny)
        #plt.imshow(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
        #plt.show()
  prs.save(f"{title}.pptx")
  print("PDF convertido!!")

if __name__ == "__main__" :
  main()